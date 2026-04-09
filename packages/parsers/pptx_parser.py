"""PPTX parser — extracts text from slides."""

from pathlib import Path

from pptx import Presentation

from packages.parsers.base import BaseParser, ParseError, ParseResult
from packages.parsers.registry import register_parser
from packages.schemas.document import FileDocument
from packages.schemas.enums import FileType
from packages.schemas.evidence import CitationAnchor, EvidenceChunk


@register_parser(FileType.PPTX)
class PptxParser(BaseParser):
    """Parse .pptx files into evidence chunks with slide-based citation anchors."""

    def parse(self, file_path: Path, file_doc: FileDocument) -> ParseResult:
        """Parse a PPTX file into evidence chunks."""
        try:
            prs = Presentation(str(file_path))
        except Exception as exc:
            raise ParseError(str(file_path), f"Cannot open PPTX: {exc}") from exc

        chunks: list[EvidenceChunk] = []
        warnings: list[str] = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            texts: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            texts.append(text)

            if texts:
                chunks.append(
                    EvidenceChunk(
                        file_id=file_doc.file_id,
                        content="\n".join(texts),
                        content_type="text",
                        citation_anchor=CitationAnchor(
                            file_id=file_doc.file_id,
                            page=slide_num,
                        ),
                    )
                )

        if not chunks:
            warnings.append("PPTX file produced no content chunks")

        file_doc.page_count = len(prs.slides)

        return ParseResult(document=file_doc, chunks=chunks, warnings=warnings)
