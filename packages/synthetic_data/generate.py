"""Generate synthetic test fixtures for all supported file formats.

Run: python -m packages.synthetic_data.generate
"""

import csv
from pathlib import Path

import docx
from openpyxl import Workbook
from pptx import Presentation
from reportlab.lib.pagesizes import A4
from reportlab.pdfgen import canvas

FIXTURES_DIR = Path(__file__).parent / "fixtures"


def _ensure_dir() -> None:
    FIXTURES_DIR.mkdir(parents=True, exist_ok=True)


def generate_clean_docx() -> Path:
    """Clean DOCX with headings, paragraphs, and a table."""
    doc = docx.Document()
    doc.add_heading("Quarterly Workforce Report", level=1)
    doc.add_paragraph("This report covers workforce metrics for Q1 2026.")
    doc.add_heading("Summary", level=2)
    doc.add_paragraph("Total headcount increased by 5% compared to Q4 2025.")
    doc.add_paragraph("Attrition rate remained stable at 3.2%.")

    table = doc.add_table(rows=4, cols=3)
    table.style = "Table Grid"
    for i, row_data in enumerate(
        [
            ("Metric", "Q4 2025", "Q1 2026"),
            ("Headcount", "1200", "1260"),
            ("New Hires", "80", "95"),
            ("Exits", "40", "35"),
        ]
    ):
        for j, val in enumerate(row_data):
            table.rows[i].cells[j].text = val

    path = FIXTURES_DIR / "clean_report.docx"
    doc.save(str(path))
    return path


def generate_messy_docx() -> Path:
    """DOCX with no headings and inconsistent formatting."""
    doc = docx.Document()
    doc.add_paragraph("workforce data 2026")
    doc.add_paragraph("")
    doc.add_paragraph("total: about 1200 maybe more")
    doc.add_paragraph("  exits unknown  ")
    doc.add_paragraph("")

    path = FIXTURES_DIR / "messy_report.docx"
    doc.save(str(path))
    return path


def generate_empty_docx() -> Path:
    """Completely empty DOCX."""
    doc = docx.Document()
    path = FIXTURES_DIR / "empty.docx"
    doc.save(str(path))
    return path


def generate_clean_pdf() -> Path:
    """Clean multi-page PDF with text."""
    path = FIXTURES_DIR / "clean_memo.pdf"
    c = canvas.Canvas(str(path), pagesize=A4)
    c.setFont("Helvetica-Bold", 16)
    c.drawString(72, 750, "Ministry Budget Memo")
    c.setFont("Helvetica", 12)
    c.drawString(72, 720, "Fiscal year 2025-2026 budget allocation summary.")
    c.drawString(72, 700, "Total allocated: SAR 450,000,000")
    c.drawString(72, 680, "Spent to date: SAR 312,500,000")
    c.drawString(72, 660, "Remaining: SAR 137,500,000")
    c.showPage()
    c.setFont("Helvetica", 12)
    c.drawString(72, 750, "Page 2: Breakdown by department.")
    c.drawString(72, 730, "HR: SAR 85,000,000")
    c.drawString(72, 710, "IT: SAR 120,000,000")
    c.drawString(72, 690, "Operations: SAR 107,500,000")
    c.showPage()
    c.save()
    return path


def generate_clean_pptx() -> Path:
    """Clean 3-slide PPTX."""
    prs = Presentation()
    layout = prs.slide_layouts[1]

    slide1 = prs.slides.add_slide(layout)
    slide1.shapes.title.text = "Q1 Performance Review"
    slide1.placeholders[1].text = "Overview of key metrics for Q1 2026"

    slide2 = prs.slides.add_slide(layout)
    slide2.shapes.title.text = "Revenue Growth"
    slide2.placeholders[1].text = "Revenue increased 12% year-over-year"

    slide3 = prs.slides.add_slide(layout)
    slide3.shapes.title.text = "Next Steps"
    slide3.placeholders[1].text = "Focus on operational efficiency in Q2"

    path = FIXTURES_DIR / "clean_slides.pptx"
    prs.save(str(path))
    return path


def generate_clean_xlsx() -> Path:
    """Clean XLSX with 2 sheets, headers, numeric data."""
    wb = Workbook()
    ws1 = wb.active
    assert ws1 is not None
    ws1.title = "Headcount"
    ws1.append(["Department", "Q4 2025", "Q1 2026"])
    ws1.append(["HR", 150, 155])
    ws1.append(["IT", 320, 340])
    ws1.append(["Operations", 730, 765])

    ws2 = wb.create_sheet("Budget")
    ws2.append(["Category", "Allocated", "Spent"])
    ws2.append(["Salaries", 200000000, 180000000])
    ws2.append(["Equipment", 50000000, 32000000])
    ws2.append(["Training", 15000000, 8000000])

    path = FIXTURES_DIR / "clean_budget.xlsx"
    wb.save(str(path))
    return path


def generate_messy_xlsx() -> Path:
    """XLSX with merged cells, missing headers, mixed types."""
    wb = Workbook()
    ws = wb.active
    assert ws is not None
    ws.title = "Data"
    ws.merge_cells("A1:C1")
    ws["A1"] = "Merged Header Row"
    ws.append([None, None, None])
    ws.append(["a", 100, "unknown"])
    ws.append([None, "N/A", 200])
    ws.append(["b", "", ""])

    path = FIXTURES_DIR / "messy_budget.xlsx"
    wb.save(str(path))
    return path


def generate_empty_xlsx() -> Path:
    """Completely empty XLSX."""
    wb = Workbook()
    path = FIXTURES_DIR / "empty.xlsx"
    wb.save(str(path))
    return path


def generate_clean_csv() -> Path:
    """Clean CSV with headers and rows."""
    path = FIXTURES_DIR / "clean_data.csv"
    with open(path, "w", newline="", encoding="utf-8") as f:
        writer = csv.writer(f)
        writer.writerow(["Name", "Department", "Salary"])
        writer.writerow(["Ahmed", "IT", "25000"])
        writer.writerow(["Sara", "HR", "22000"])
        writer.writerow(["Khalid", "Operations", "21000"])
        writer.writerow(["Fatima", "IT", "26000"])
    return path


def generate_empty_csv() -> Path:
    """Empty CSV."""
    path = FIXTURES_DIR / "empty.csv"
    path.write_text("", encoding="utf-8")
    return path


def main() -> None:
    """Generate all synthetic fixtures."""
    _ensure_dir()

    generators = [
        generate_clean_docx,
        generate_messy_docx,
        generate_empty_docx,
        generate_clean_pdf,
        generate_clean_pptx,
        generate_clean_xlsx,
        generate_messy_xlsx,
        generate_empty_xlsx,
        generate_clean_csv,
        generate_empty_csv,
    ]

    for gen in generators:
        path = gen()
        print(f"  created: {path.name}")  # noqa: T201

    print(f"\n{len(generators)} fixtures generated in {FIXTURES_DIR}")  # noqa: T201


if __name__ == "__main__":
    main()
